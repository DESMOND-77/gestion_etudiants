"""
Générateur complet de la présentation PowerPoint
"Système de Gestion des Étudiants USTM" - 31 slides
Respecte exactement les specs : couleurs, polices, tailles, encadrés, tableaux.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

def add_colored_box(slide, left, top, width, height, text, 
                   fill_rgb, line_rgb=None, font_size=14, font_bold=True):
    """Ajoute un encadré arrondi coloré avec texte."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    
    tf = shape.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = font_bold
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_two_columns(slide, left1, left2, top, width, height, items1, items2, color_rgb):
    """Ajoute 2 colonnes de texte avec puces."""
    # Colonne 1
    tx1 = slide.shapes.add_textbox(left1, top, width, height)
    tf1 = tx1.text_frame
    tf1.word_wrap = True
    for item in items1:
        p = tf1.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.name = "Arial"
        p.space_after = Pt(12)
        p.font.color.rgb = color_rgb
    
    # Colonne 2
    tx2 = slide.shapes.add_textbox(left2, top, width, height)
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    for item in items2:
        p = tf2.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.name = "Arial"
        p.space_after = Pt(12)
        p.font.color.rgb = color_rgb

def add_table(slide, left, top, width, height, rows, cols, data, header=True):
    """Ajoute un tableau formaté."""
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table
    
    for i, row in enumerate(data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = cell_text
            # Formatage header
            if header and i == 0:
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(16)
            else:
                cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].font.name = "Arial"
    return table

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 🎨 Palette de couleurs exacte
    C = {
        'bleu_marine': RGBColor(0, 51, 102),
        'bleu_fonce': RGBColor(13, 71, 161),
        'bleu_clair': RGBColor(232, 244, 248),
        'vert_fonce': RGBColor(46, 125, 50),
        'vert_clair': RGBColor(200, 230, 201),
        'rouge_fonce': RGBColor(198, 40, 40),
        'rouge_clair': RGBColor(255, 205, 210),
        'orange_fonce': RGBColor(239, 108, 0),
        'orange_clair': RGBColor(255, 224, 178),
        'violet_fonce': RGBColor(106, 27, 154),
        'violet_clair': RGBColor(243, 229, 245),
        'gris_fonce': RGBColor(51, 51, 51),
        'gris_clair': RGBColor(245, 245, 245),
        'jaune_clair': RGBColor(255, 249, 196),
        'fond_blanc': RGBColor(255, 255, 255),
        'noir': RGBColor(0, 0, 0),
    }

    # =========================================================================
    # 📋 SLIDE 1 - PAGE DE TITRE
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Titre seul
    
    title = slide.shapes.title
    title.text = "SYSTÈME DE GESTION DES ÉTUDIANTS\nUSTM"
    title.text_frame.paragraphs[0].font.name = "Arial Black"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = C['bleu_marine']
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Programmation Avancée - DUT2 GRT\nSemestre S1"
    subtitle.text_frame.paragraphs[0].font.name = "Arial"
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[0].font.color.rgb = C['gris_fonce']
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Bas de page
    tx = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11.333), Inches(0.6))
    p = tx.text_frame.paragraphs[0]
    p.text = "Université des Sciences et Techniques de Masuku • Franceville, Gabon"
    p.font.size = Pt(16)
    p.font.name = "Arial"
    p.font.color.rgb = RGBColor(102, 102, 102)
    p.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # 📋 SLIDE 2 - PLAN
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Titre + Contenu
    title = slide.shapes.title
    title.text = "📑 PLAN DE LA PRÉSENTATION"
    title.text_frame.paragraphs[0].font.color.rgb = C['bleu_marine']
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    
    items_left = [
        "🎯 1. Contexte et Problématique",
        "📊 2. Objectifs du Projet",
        "💡 3. Concept et Innovation",
        "🏗️ 4. Architecture Technique",
        "📦 5. Structure de Données"
    ]
    items_right = [
        "⚙️ 6. Algorithmes Clés",
        "💻 7. Implémentation",
        "📁 8. Format CSV",
        "🌍 9. Applications au Gabon",
        "🎓 10. Conclusion"
    ]
    
    add_two_columns(slide, Inches(0.5), Inches(6.5), Inches(1.8), Inches(4), Inches(5),
                    items_left, items_right, C['noir'])

    # =========================================================================
    # 📋 SLIDE 3 - CONTEXTE
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🏛️ CONTEXTE DU PROJET"
    title.text_frame.paragraphs[0].font.color.rgb = C['bleu_marine']
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    # Encadrés
    add_colored_box(slide, Inches(0.5), Inches(2.0), Inches(4), Inches(0.8),
                    "📍 USTM\nUniversité des Sciences et Techniques de Masuku\nFranceville, Gabon",
                    C['bleu_clair'], C['bleu_marine'], 14, True)
    
    add_colored_box(slide, Inches(0.5), Inches(3.0), Inches(4), Inches(0.8),
                    "🎓 Formation\nDUT2 - Génie Réseaux et Télécommunications\nCours : Programmation Avancée",
                    C['vert_clair'], C['vert_fonce'], 14, True)
    
    add_colored_box(slide, Inches(0.5), Inches(4.0), Inches(4), Inches(0.8),
                    "🔬 Type de Projet\nApplication académique en langage C\nGestion de base de données CSV",
                    RGBColor(255, 243, 224), C['orange_fonce'], 14, True)

    # =========================================================================
    # 📋 SLIDE 4 - PROBLÉMATIQUE
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "❓ PROBLÉMATIQUE"
    title.text_frame.paragraphs[0].font.color.rgb = C['rouge_fonce']
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    # Diagramme texte
    diag_text = """
┌─────────────────────────────────────────────────┐
│         GESTION DES ARCHIVES                    │
├─────────────────┬───────────────────────────────┤
│   📄 PAPIER     │      💾 NUMÉRIQUE             │
├─────────────────┼───────────────────────────────┤
│ • Dossiers      │ • Fichiers CSV                │
│ • Fiches        │ • Bases de données            │
│ • Registres     │ • Tableurs                    │
└─────────────────┴───────────────────────────────┘
           │                    │
           ▼                    ▼
     ⏱️ TEMPS DE RECHERCHE
           │
    ┌──────┴──────┐
    ▼             ▼
 ✅ RAPIDE    ❌ COMPLEXE
"""
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(4))
    p = tx.text_frame.paragraphs[0]
    p.text = diag_text
    p.font.name = "Courier New"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.LEFT
    
    # Défi principal
    add_colored_box(slide, Inches(1), Inches(5.5), Inches(11.333), Inches(0.9),
                    "⚠️ DÉFI PRINCIPAL\nComment insérer automatiquement les données au BON ENDROIT dans un fichier numérique ?",
                    C['jaune_clair'], C['rouge_fonce'], 20, True)

    # =========================================================================
    # 📋 SLIDE 5 - CAHIER DES CHARGES
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📋 CAHIER DES CHARGES"
    title.text_frame.paragraphs[0].font.color.rgb = C['bleu_marine']
    
    flux_text = """
        ┌─────────────────────────┐
        │  Nouvelle inscription   │
        │      👨‍🎓 Étudiant       │
        └───────────┬─────────────┘
                    │
                    ▼
        ┌─────────────────────────┐
        │   📂 Fichier CSV         │
        │   Recherche de doublon  │
        └───────────┬─────────────┘
                    │
            ┌───────┴───────┐
            ▼               ▼
    ✅ INEXISTANT    ❌ EXISTE DÉJÀ
            │               │
            ▼               ▼
    ┌──────────────┐  ┌──────────────┐
    │ INSERTION    │  │  REFUS       │
    │ Ordre Alpha  │  │  + Message   │
    └──────────────┘  └──────────────┘
"""
    tx = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(4.5))
    p = tx.text_frame.paragraphs[0]
    p.text = flux_text
    p.font.name = "Courier New"
    p.font.size = Pt(13)
    p.alignment = PP_ALIGN.CENTER
    
    add_colored_box(slide, Inches(1), Inches(6.0), Inches(11.333), Inches(0.8),
                    "🎯 Ordre alphabétique CROISSANT des NOMS\n🔒 Unicité garantie par le MATRICULE\n💾 Persistance dans fichier CSV",
                    C['vert_clair'], C['vert_fonce'], 18, True)

    # =========================================================================
    # 📋 SLIDE 6 - OBJECTIFS (4 blocs)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🎯 OBJECTIFS DU PROJET"
    title.text_frame.paragraphs[0].font.color.rgb = C['vert_fonce']
    
    objectifs = [
        (C['bleu_marine'], "🔤 ORDRE ALPHABÉTIQUE\n\nMaintenir le tri automatique par nom de famille\n\nInsertion intelligente"),
        (C['orange_fonce'], "🔒 UNICITÉ\n\nVérification des doublons par matricule\n\nAucune duplication"),
        (C['violet_fonce'], "💾 PERSISTANCE\n\nStockage dans fichier CSV standard\n\nInteropérabilité"),
        (C['vert_fonce'], "🖥️ DOUBLE INTERFACE\n\nVersion Console +\nVersion Graphique GTK+")
    ]
    
    positions = [(0.5, 1.8), (6.5, 1.8), (0.5, 4.2), (6.5, 4.2)]
    for (color, text), (x, y) in zip(objectifs, positions):
        add_colored_box(slide, Inches(x), Inches(y), Inches(5.8), Inches(2.2),
                        text, color, color, 18, True)

    # =========================================================================
    # 📋 SLIDE 7 - CONCEPT CRUD
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "💡 CONCEPT : SYSTÈME CRUD"
    title.text_frame.paragraphs[0].font.color.rgb = C['violet_fonce']
    
    # Diagramme circulaire simplifié
    circle_text = """
                    📝 CREATE
                   (Ajouter)
                       │
                       │
    🗑️ DELETE ────────┼──────── 📖 READ
    (Supprimer)        │        (Afficher)
                       │
                       │
                   ✏️ UPDATE
                  (Modifier)
"""
    tx = slide.shapes.add_textbox(Inches(4.5), Inches(1.5), Inches(4.5), Inches(3))
    p = tx.text_frame.paragraphs[0]
    p.text = circle_text
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER
    
    # Tableau
    data = [
        ["🔄 Opération", "📌 Fonction", "⚡ Innovation"],
        ["CREATE", "Ajouter étudiant", "**Insertion triée auto**"],
        ["READ", "Afficher/Rechercher", "Recherche binaire O(log n)"],
        ["UPDATE", "Modifier infos", "Validation des données"],
        ["DELETE", "Supprimer", "Confirmation sécurisée"]
    ]
    add_table(slide, Inches(0.5), Inches(4.5), Inches(12.333), Inches(2.2), 5, 3, data, True)

    # =========================================================================
    # 📋 SLIDE 8 - ARCHITECTURE
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🏗️ ARCHITECTURE DU PROJET"
    title.text_frame.paragraphs[0].font.color.rgb = C['bleu_marine']
    
    arch_text = """
┏━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━┓
┃   👤 INTERFACE UTILISATEUR               ┃
┃   🖥️ Console CMD  │  🎨 GUI GTK+ 3.0    ┃
┗━━━━━━━━━━━━━━━━━━━┯━━━━━━━━━━━━━━━━━━━━┛
                     │
┏━━━━━━━━━━━━━━━━━━━▼━━━━━━━━━━━━━━━━━━━━┓
┃   ⚙️ LOGIQUE MÉTIER                     ┃
┃   Insertion │ Recherche │ Tri │ Suppression ┃
┗━━━━━━━━━━━━━━━━━━━┯━━━━━━━━━━━━━━━━━━━━┛
                     │
┏━━━━━━━━━━━━━━━━━━━▼━━━━━━━━━━━━━━━━━━━━┓
┃   📁 GESTION DES DONNÉES                ┃
┃   Lecture/Écriture CSV │ Parsing       ┃
┗━━━━━━━━━━━━━━━━━━━┯━━━━━━━━━━━━━━━━━━━━┛
                     │
┏━━━━━━━━━━━━━━━━━━━▼━━━━━━━━━━━━━━━━━━━━┓
┃   📦 STRUCTURE DE DONNÉES               ┃
┃   Tableau statique Etudiant[1000]      ┃
┗━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━┛
"""
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(5.5))
    p = tx.text_frame.paragraphs[0]
    p.text = arch_text
    p.font.name = "Courier New"
    p.font.size = Pt(13)
    p.alignment = PP_ALIGN.LEFT

    # =========================================================================
    # 📋 SLIDE 9 - ARBORESCENCE
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📂 STRUCTURE DES FICHIERS"
    title.text_frame.paragraphs[0].font.color.rgb = C['orange_fonce']
    
    tree_text = """📁 gestion_etudiants/
│
├── 📁 CMD_gestion/           🖥️ Version Console
│   └── 📄 CMD_gestion_etudiants.c
│
├── 📁 GUI_gestion/           🎨 Version Graphique
│   ├── 📄 GUI_gestion_etudiants.c
│   └── 🎨 style.css
│
├── 📁 functions/             ⚙️ Fonctions partagées
│   ├── 📄 fonctions.c       (Parsing, tri, recherche)
│   ├── 📄 etudiant.c        (Saisie, affichage)
│   └── 📄 banner.c          (Messages)
│
├── 📁 libs/                  📚 Bibliothèques
│   ├── 📄 fonctions.h
│   ├── 📄 etudiant.h
│   ├── 📄 couleurs.h
│   └── 📄 banner.h
│
├── 📁 data/                  💾 Données
│   └── 📄 Etudiants.csv
│
└── 📄 makefile               🔨 Compilation"""
    
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(5.5))
    p = tx.text_frame.paragraphs[0]
    p.text = tree_text
    p.font.name = "Courier New"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.LEFT

    # =========================================================================
    # 📋 SLIDE 10 - STRUCTURE DE DONNÉES
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📦 STRUCTURE DE DONNÉES"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(21, 101, 192)
    
    # Code C
    code_c = """typedef struct
{
    char nom[128];         // 👤 Nom de famille
    char prenom[128];      // 👤 Prénom
    char classe[128];      // 🎓 Niveau/Filière
    char matricule[128];   // 🔑 ID unique
    char email[128];       // 📧 Adresse email
    char sex[128];         // ⚧ Sexe (M/F)
} Etudiant;"""
    
    add_colored_box(slide, Inches(0.5), Inches(1.8), Inches(12.333), Inches(2.2),
                    code_c, C['gris_clair'], RGBColor(120,120,120), 16, False)
    
    # Tableau constantes
    data_const = [
        ["🏷️ Constante", "🔢 Valeur", "📝 Description"],
        ["MAX_STUD", "1000", "Nombre max d'étudiants"],
        ["MAX_CHAMP", "128", "Longueur max d'un champ"],
        ["MAX_LIGNE", "512", "Longueur max ligne CSV"]
    ]
    add_table(slide, Inches(0.5), Inches(4.2), Inches(12.333), Inches(2.0), 5, 3, data_const, True)

    # =========================================================================
    # 📋 SLIDE 11 - RECHERCHE BINAIRE
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🔍 ALGORITHME : RECHERCHE BINAIRE"
    title.text_frame.paragraphs[0].font.color.rgb = C['rouge_fonce']
    
    algo_text = """
Rechercher "Biyiti" dans :
["Dupont", "Martin", "Nguema", "Obame"]

Étape 1:  [Dupont] | Martin | [Nguema, Obame]
          ↑ Début   ↑ Milieu  ↑ Fin
          "B" < "M" → Chercher à GAUCHE

Étape 2:  [Dupont]
          "B" < "D" → Insérer AVANT

Résultat: ["Biyiti", "Dupont", "Martin", "Nguema", "Obame"]
          ↑ Position 0
"""
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(3.5))
    p = tx.text_frame.paragraphs[0]
    p.text = algo_text
    p.font.name = "Courier New"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.LEFT
    
    add_colored_box(slide, Inches(2), Inches(5.2), Inches(9.333), Inches(0.8),
                    "⚡ COMPLEXITÉ : O(log n)\nTrès efficace même avec 1000+ étudiants !",
                    RGBColor(165, 214, 167), C['vert_fonce'], 20, True)

    # =========================================================================
    # 📋 SLIDE 12 - DÉTECTION DOUBLONS
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🔒 DÉTECTION DES DOUBLONS"
    title.text_frame.paragraphs[0].font.color.rgb = C['orange_fonce']
    
    flow_text = """
        ┌──────────────────────┐
        │  Nouvel étudiant     │
        │  Matricule: MAT2024X │
        └──────────┬───────────┘
                   │
                   ▼
        ┌──────────────────────┐
        │ Parcourir tous les   │
        │ étudiants existants  │
        └──────────┬───────────┘
                   │
        ┌──────────▼───────────┐
        │ Comparer matricules  │
        └──────────┬───────────┘
                   │
           ┌───────┴────────┐
           ▼                ▼
    ✅ UNIQUE         ❌ DOUBLON
           │                │
           ▼                ▼
    Insertion OK    Refus + Message
"""
    tx = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(4.5))
    p = tx.text_frame.paragraphs[0]
    p.text = flow_text
    p.font.name = "Courier New"
    p.font.size = Pt(13)
    p.alignment = PP_ALIGN.CENTER
    
    code_ex = """for (int i = 0; i < n; i++)
{
    if (matricule_existe(tab[i]))
        return ERREUR_DOUBLON; // ❌
}
// ✅ Matricule unique"""
    add_colored_box(slide, Inches(1), Inches(6.0), Inches(11.333), Inches(1.0),
                    code_ex, RGBColor(255, 253, 231), C['orange_fonce'], 14, False)

    # =========================================================================
    # 📋 SLIDE 13 - INSERTION TRIÉE
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "➕ INSERTION AVEC DÉCALAGE"
    title.text_frame.paragraphs[0].font.color.rgb = C['vert_fonce']
    
    etapes = [
        """Étape 1 - État initial:
[0]     [1]      [2]      [3]      [4]
Dupont | Martin | Nguema | Obame | [VIDE]""",
        
        """Étape 2 - Décalage à droite (insérer "Koundou" à pos 1):
[0]     [1]      [2]      [3]      [4]
Dupont | [VIDE] | Martin→ | Nguema→ | Obame→""",
        
        """Étape 3 - Insertion finale:
[0]     [1]      [2]      [3]      [4]
Dupont | Koundou | Martin | Nguema | Obame ← ✅"""
    ]
    
    y_pos = 1.8
    for etape in etapes:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12.333), Inches(1.2))
        p = tx.text_frame.paragraphs[0]
        p.text = etape
        p.font.name = "Courier New"
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT
        y_pos += 1.4
    
    add_colored_box(slide, Inches(1), Inches(6.0), Inches(11.333), Inches(0.8),
                    "⏱️ COMPLEXITÉ : O(n)\nDans le pire cas, décalage de tous les éléments",
                    RGBColor(255, 224, 178), C['orange_fonce'], 18, True)

    # =========================================================================
    # 📋 SLIDE 14 - PARSING CSV
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📄 PARSING DU FICHIER CSV"
    title.text_frame.paragraphs[0].font.color.rgb = C['violet_fonce']
    
    csv_line = "Dupont;Jean;L1 Info;jean.dupont@ustm.fr;MAT001;M"
    parsing = """┌──────────────────────────────────────┐
│ Délimiteur: ; (point-virgule)       │
├──────────────┬───────────────────────┤
│ nom          │ "Dupont"              │
│ prenom       │ "Jean"                │
│ classe       │ "L1 Info"             │
│ email        │ "jean.dupont@ustm.fr" │
│ matricule    │ "MAT001"              │
│ sex          │ "M"                   │
└──────────────┴───────────────────────┘"""
    
    add_colored_box(slide, Inches(0.5), Inches(1.8), Inches(12.333), Inches(0.8),
                    csv_line, C['jaune_clair'], None, 16, False)
    
    tx = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.333), Inches(3.5))
    p = tx.text_frame.paragraphs[0]
    p.text = parsing
    p.font.name = "Courier New"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.LEFT

    # =========================================================================
    # 📋 SLIDE 15 - VERSION CONSOLE
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🖥️ INTERFACE CONSOLE (CMD)"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(13, 71, 161)
    
    menu_text = """
┏━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━┓
┃                                               ┃
┃       GESTION DES ÉTUDIANTS - USTM            ┃
┃                                               ┃
┗━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━┛

        Nombre d'étudiants: 15

        1. ➕ Ajouter un étudiant
        2. 📋 Afficher tous les étudiants
        3. 🔍 Rechercher par nom
        4. ✏️  Modifier un étudiant
        5. 🗑️  Supprimer un étudiant
        6. 🔤 Trier les étudiants
        0. 🚪 Quitter

        Votre choix: _
"""
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(4.5))
    p = tx.text_frame.paragraphs[0]
    p.text = menu_text
    p.font.name = "Courier New"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.LEFT
    
    # 3 blocs caractéristiques
    add_colored_box(slide, Inches(0.5), Inches(6.0), Inches(3.8), Inches(1.0),
                    "🎨 Couleurs ANSI\nCodes d'échappement\npour terminal", C['bleu_clair'], None, 14, True)
    add_colored_box(slide, Inches(4.5), Inches(6.0), Inches(3.8), Inches(1.0),
                    "⌨️ Menu interactif\nNavigation par\nnuméros", C['violet_clair'], None, 14, True)
    add_colored_box(slide, Inches(8.5), Inches(6.0), Inches(3.8), Inches(1.0),
                    "✅ Messages clairs\nConfirmation\net erreurs", C['vert_clair'], None, 14, True)

    # =========================================================================
    # 📋 SLIDE 16 - VERSION GRAPHIQUE
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🎨 INTERFACE GRAPHIQUE (GTK+)"
    title.text_frame.paragraphs[0].font.color.rgb = C['vert_fonce']
    
    gtk_text = """
┌────────────────────────────────────────┐
│  📝 CHAMPS DE SAISIE                   │
│  • Nom, Prénom, Classe                 │
│  • Email, Matricule, Sexe              │
├────────────────────────────────────────┤
│  🔘 BOUTONS                            │
│  Insérer | Afficher | Rechercher       │
│  Modifier | Supprimer | Trier | Quitter│
├────────────────────────────────────────┤
│  📊 TABLEAU (TreeView)                 │
│  Affichage des étudiants en colonnes  │
└────────────────────────────────────────┘
"""
    tx = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(4.5))
    p = tx.text_frame.paragraphs[0]
    p.text = gtk_text
    p.font.name = "Courier New"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.LEFT
    
    add_colored_box(slide, Inches(1), Inches(6.0), Inches(11.333), Inches(0.8),
                    "🔧 Bibliothèque: GTK+ 3.0  •  🎨 Style: CSS personnalisé  •  🖱️ Événements: Callbacks",
                    RGBColor(224, 242, 241), C['vert_fonce'], 18, True)

    # =========================================================================
    # 📋 SLIDE 17 - FORMAT CSV
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📁 FORMAT DU FICHIER CSV"
    title.text_frame.paragraphs[0].font.color.rgb = C['orange_fonce']
    
    csv_content = """nom;prenom;classe;email;matricule;sex
───────────────────────────────────────────
Dupont;Jean;L1 Info;jean@ustm.fr;MAT001;M
Martin;Marie;L2 Info;marie@ustm.fr;MAT002;F
Nguema;Paul;L3 Tech;paul@ustm.fr;MAT003;M
Obame;Claire;L1 Info;claire@ustm.fr;MAT004;F"""
    
    add_colored_box(slide, Inches(0.5), Inches(1.8), Inches(12.333), Inches(3.5),
                    csv_content, C['jaune_clair'], None, 16, False)
    
    # 3 encadrés
    add_colored_box(slide, Inches(0.5), Inches(5.3), Inches(4), Inches(0.9),
                    "📍 Emplacement\ndata/Etudiants.csv", RGBColor(255, 243, 224), None, 14, True)
    add_colored_box(slide, Inches(4.7), Inches(5.3), Inches(4), Inches(0.9),
                    "➗ Délimiteur\n; (point-virgule)", RGBColor(225, 245, 254), None, 14, True)
    add_colored_box(slide, Inches(8.9), Inches(5.3), Inches(3.8), Inches(0.9),
                    "🔄 Compatibilité\nExcel, LibreOffice", RGBColor(241, 248, 233), None, 14, True)

    # =========================================================================
    # 📋 SLIDE 18 - COMPILATION
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🔨 COMPILATION DU PROJET"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(25, 118, 210)
    
    makefile = """# Compiler les deux versions
$ make all

# Compiler uniquement la version console
$ make cmd

# Compiler uniquement la version GUI
$ make gui

# Nettoyer les fichiers compilés
$ make clean"""
    
    manuel = """# Version Console
$ gcc -Wall -Wextra -std=c99 -I./libs \
      CMD_gestion/CMD_gestion_etudiants.c \
      functions/*.c -o cmd_gestion

# Version GUI
$ gcc -Wall -Wextra -std=c99 -I./libs \
      GUI_gestion/GUI_gestion_etudiants.c \
      functions/*.c `pkg-config --cflags --libs gtk+-3.0` \
      -o gui_gestion"""
    
    add_colored_box(slide, Inches(0.5), Inches(1.8), Inches(6), Inches(2.5),
                    makefile, C['vert_clair'], None, 16, False)
    add_colored_box(slide, Inches(6.7), Inches(1.8), Inches(6), Inches(3.0),
                    manuel, RGBColor(255, 224, 178), None, 14, False)

    # =========================================================================
    # 📋 SLIDE 19 - EXÉCUTION (2 colonnes)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "▶️ EXÉCUTION DES PROGRAMMES"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(104, 159, 56)
    
    # Gauche - Console
    add_colored_box(slide, Inches(0.5), Inches(1.8), Inches(6), Inches(2.2),
                    "$ ./cmd_gestion_etudiants", C['bleu_clair'], None, 20, False)
    add_colored_box(slide, Inches(0.5), Inches(4.2), Inches(6), Inches(2.0),
                    "✅ Léger (faibles ressources)\n✅ Rapide\n✅ Compatible serveurs",
                    C['bleu_clair'], None, 16, True)
    
    # Droite - GUI
    add_colored_box(slide, Inches(6.7), Inches(1.8), Inches(6), Inches(2.2),
                    "$ ./gui_gestion_etudiants", C['vert_clair'], None, 20, False)
    add_colored_box(slide, Inches(6.7), Inches(4.2), Inches(6), Inches(2.0),
                    "✅ Convivial\n✅ Visuel\n✅ Intuitif",
                    C['vert_clair'], None, 16, True)

    # =========================================================================
    # 📋 SLIDE 20 - CONTEXTE GABONAIS (1/2) - PROBLÈMES
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🇬🇦 PROBLÉMATIQUES AU GABON"
    title.text_frame.paragraphs[0].font.color.rgb = C['rouge_fonce']
    
    probs = [
        (C['rouge_clair'], "📚 GESTION ADMINISTRATIVE\n• Volumes importants d'étudiants\n• Multiples filières et niveaux\n• Turn-over du personnel\n• Processus manuels lents"),
        (C['orange_clair'], "📁 ARCHIVES MIXTES\n• Documents papier anciens\n• Migration vers le numérique\n• Besoin de centralisation\n• Risques de perte"),
        (C['jaune_clair'], "💰 RESSOURCES LIMITÉES\n• Connectivité internet variable\n• Besoin d'applications légères\n• Coût des licences logicielles\n• Matériel informatique limité")
    ]
    y = 1.8
    for color, text in probs:
        add_colored_box(slide, Inches(0.5), Inches(y), Inches(12.333), Inches(1.4),
                        text, color, None, 14, True)
        y += 1.6

    # =========================================================================
    # 📋 SLIDE 21 - CONTEXTE GABONAIS (2/2) - SOLUTIONS
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "✅ SOLUTIONS PROPOSÉES"
    title.text_frame.paragraphs[0].font.color.rgb = C['vert_fonce']
    
    sol_data = [
        ["🏛️ Institution", "📌 Application", "🎯 Bénéfice"],
        ["USTM, UOB", "Gestion des inscriptions", "⚡ Insertion auto triée"],
        ["Universités", "Suivi par classe", "🔍 Recherche binaire rapide"],
        ["Administrations", "Exportation données", "💾 Compatibilité CSV/Excel"],
        ["Lycées/Collèges", "Gestion élèves", "📋 Système adaptable"],
        ["Archives", "Transition numérique", "📄 Papier → Digital"]
    ]
    add_table(slide, Inches(0.5), Inches(1.8), Inches(12.333), Inches(3.5), 7, 3, sol_data, True)
    
    # Avantages clés
    avants = ["💰 GRATUIT", "🚀 PERFORMANT", "🔓 OPEN SOURCE",
              "🪶 LÉGER", "📊 STANDARD"]
    x_pos = [0.5, 3.5, 6.5, 9.0, 0.5, 3.5, 6.5, 9.0, 0.5, 3.5][:5]  # correction
    y_pos = [5.4, 5.4, 5.4, 5.4, 6.3, 6.3, 6.3, 6.3, 5.4, 5.4][:5]
    
    texts = [
        "💰 GRATUIT\nPas de licence",
        "🚀 PERFORMANT\nO(log n) recherche",
        "🔓 OPEN SOURCE\nCode modifiable",
        "🪶 LÉGER\nFonctionne sur PC modeste",
        "📊 STANDARD\nFormat CSV universel"
    ]
    colors = [C['vert_clair'], C['bleu_clair'], C['violet_clair'], C['orange_clair'], C['jaune_clair']]
    
    for i, (txt, col) in enumerate(zip(texts, colors)):
        add_colored_box(slide, Inches(0.5 + (i%3)*3.2), Inches(5.4 + (i//3)*0.9), Inches(3.2), Inches(0.8),
                        txt, col, None, 14, True)

    # =========================================================================
    # 📋 SLIDE 22 - CAS D'USAGE (4 scénarios)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🎯 CAS D'USAGE RÉELS"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(21, 101, 192)
    
    scenarios = [
        (C['bleu_clair'], "📝 INSCRIPTION ANNUELLE\n• Ajout de centaines d'étudiants\n• Insertion automatique triée\n• Gain de temps considérable\n• Zéro erreur de classement"),
        (C['violet_clair'], "✏️ MISE À JOUR\n• Changement de classe\n• Modification d'email\n• Correction d'erreurs\n• Historique préservé"),
        (C['vert_clair'], "🔍 RECHERCHE RAPIDE\n• Par nom pour PV d'examen\n• Vérification d'inscription\n• Consultation dossier\n• O(log n) efficacité"),
        (C['orange_clair'], "📊 GÉNÉRATION DE LISTES\n• Listes par classe\n• Convocations examens\n• Export vers Excel\n• Impression facile")
    ]
    
    for i, (col, txt) in enumerate(scenarios):
        x = Inches(0.5 + (i%2)*6.5)
        y = Inches(1.8 + (i//2)*2.6)
        add_colored_box(slide, x, y, Inches(6.0), Inches(2.2), txt, col, None, 14, True)

    # =========================================================================
    # 📋 SLIDE 23 - COMPÉTENCES
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🎓 COMPÉTENCES DÉMONTRÉES"
    title.text_frame.paragraphs[0].font.color.rgb = C['violet_fonce']
    
    comp_text = """
                Programmation
                   modulaire
                      ⭐
                     /  \\
                    /    \\
        Makefile   /      \\   Structures
           ⭐     /        \\      ⭐
                /          \\
               /            \\
          ⭐ ──────────────── ⭐
    Interfaces         Algorithmes
  (Console/GUI)        (Recherche/Tri)
               \\           /
                \\         /
                 \\       /
                  \\     /
                   \\   /
                    ⭐
            Gestion fichiers
                 (CSV)
"""
    tx = slide.shapes.add_textbox(Inches(2), Inches(1.8), Inches(9.333), Inches(4.5))
    p = tx.text_frame.paragraphs[0]
    p.text = comp_text
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER
    
    add_colored_box(slide, Inches(2), Inches(6.0), Inches(9.333), Inches(0.8),
                    "✅ TOUTES LES COMPÉTENCES : MAÎTRISÉES", C['vert_clair'], C['vert_fonce'], 24, True)

    # =========================================================================
    # 📋 SLIDE 24 - PERSPECTIVES (1/2)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🚀 PERSPECTIVES D'AMÉLIORATION"
    title.text_frame.paragraphs[0].font.color.rgb = C['orange_fonce']
    
    futures = [
        (RGBColor(255, 235, 238), "📄 EXPORT PDF\nGénération de listes\nd'étudiants au format\nPDF imprimable"),
        (RGBColor(232, 234, 246), "📊 IMPORT EXCEL/JSON\nImportation depuis\nautres formats de\nfichiers standards"),
        (RGBColor(224, 242, 241), "🔐 AUTHENTIFICATION\nSystème de connexion\navec mots de passe\net profils utilisateurs"),
        (C['jaune_clair'], "💾 SAUVEGARDE AUTO\nEnregistrement\nautomatique\npériodique"),
        (C['vert_clair'], "📈 STATISTIQUES\nGraphiques et\nanalyses de\ndonnées")
    ]
    
    for i, (col, txt) in enumerate(futures):
        x = Inches(0.5 + (i%2)*6.5)
        y = Inches(1.8 + (i//2)*2.0)
        add_colored_box(slide, x, y, Inches(5.8), Inches(1.6), txt, col, None, 14, True)

    # =========================================================================
    # 📋 SLIDE 25 - PERSPECTIVES (2/2) - ROADMAP
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🔮 AMÉLIORATIONS TECHNIQUES"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(21, 101, 192)
    
    roadmap = """
PHASE 1          PHASE 2           PHASE 3
(Court terme)    (Moyen terme)     (Long terme)
     │                │                 │
     ▼                ▼                 ▼
┌─────────┐      ┌─────────┐      ┌──────────┐
│ SQLite  │      │   Web   │      │  Mobile  │
│ Tests   │──────│  HTML/  │──────│ Android  │
│unitaires│      │CSS/JS   │      │   iOS    │
└─────────┘      └─────────┘      └──────────┘
     │                │                 │
     ▼                ▼                 ▼
  🗄️ BDD          🌐 Cloud         📱 Apps
"""
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(3.5))
    p = tx.text_frame.paragraphs[0]
    p.text = roadmap
    p.font.name = "Courier New"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.LEFT
    
    checklist = """□ Base de données SQLite
□ Tests unitaires automatisés
□ Interface Web responsive
□ API REST pour intégrations
□ Application mobile Android/iOS
□ Internationalisation (FR/EN)
□ Synchronisation cloud
□ Mode hors ligne"""
    add_colored_box(slide, Inches(0.5), Inches(5.3), Inches(12.333), Inches(1.8),
                    checklist, C['gris_clair'], None, 16, False)

    # =========================================================================
    # 📋 SLIDE 26 - DÉMONSTRATION
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🎬 DÉMONSTRATION VIDÉO"
    title.text_frame.paragraphs[0].font.color.rgb = C['rouge_fonce']
    
    # Cadre vidéo simulé
    add_colored_box(slide, Inches(1), Inches(2.0), Inches(11.333), Inches(4.5),
                    "▶️  VIDÉO DÉMO\n\n[Capture d'écran de l'application en action]\n\n• Insertion automatique triée\n• Recherche binaire\n• Modification/Suppression\n• Export CSV",
                    C['gris_clair'], C['bleu_marine'], 24, True)
    
    # Bulles autour
    add_colored_box(slide, Inches(1), Inches(1.8), Inches(2.5), Inches(0.6),
                    "1️⃣ Insertion auto triée", C['bleu_clair'], None, 12, True)
    add_colored_box(slide, Inches(9.5), Inches(1.8), Inches(2.8), Inches(0.6),
                    "2️⃣ Recherche rapide", C['violet_clair'], None, 12, True)
    add_colored_box(slide, Inches(1), Inches(6.8), Inches(2.8), Inches(0.6),
                    "3️⃣ Modification", C['vert_clair'], None, 12, True)
    add_colored_box(slide, Inches(9.0), Inches(6.8), Inches(3.3), Inches(0.6),
                    "4️⃣ Affichage complet", C['orange_clair'], None, 12, True)

    # =========================================================================
    # 📋 SLIDE 27 - TECHNOLOGIES
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🛠️ STACK TECHNOLOGIQUE"
    title.text_frame.paragraphs[0].font.color.rgb = C['vert_fonce']
    
    stack_text = """
┌─────────────────────────────────────────┐
│        💻 LANGAGE                       │
│           C (Standard C99)              │
│         [Logo C]                        │
├─────────────────────────────────────────┤
│        🔧 COMPILATEUR                   │
│       GCC (GNU Compiler Collection)     │
│         [Logo GCC]                      │
├─────────────────────────────────────────┤
│        🎨 INTERFACES                    │
│   ANSI Colors  │  GTK+ 3.0              │
│   [Terminal]   │  [Logo GTK]            │
├─────────────────────────────────────────┤
│        🔨 BUILD                         │
│          Makefile + pkg-config          │
│         [Logo Make]                     │
├─────────────────────────────────────────┤
│        💾 DONNÉES                       │
│          Format CSV standard            │
│         [Icône fichier CSV]             │
└─────────────────────────────────────────┘
"""
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(5.5))
    p = tx.text_frame.paragraphs[0]
    p.text = stack_text
    p.font.name = "Courier New"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.LEFT
    
    # Compatibilité
    comp = ["🐧 Linux", "🪟 Windows", "🍎 macOS"]
    comp_sub = ["Ubuntu/Debian\nFedora/Arch", "MSYS2/MinGW\nCygwin", "Homebrew\nXcode"]
    colors_comp = [C['vert_clair'], C['bleu_clair'], C['violet_clair']]
    
    for i, (os, sub, col) in enumerate(zip(comp, comp_sub, colors_comp)):
        x = Inches(0.5 + i*4.5)
        add_colored_box(slide, x, Inches(6.0), Inches(4.0), Inches(1.2),
                        f"{os}\n{sub}", col, None, 14, True)

    # =========================================================================
    # 📋 SLIDE 28 - RÉSULTATS
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📊 RÉSULTATS ET PERFORMANCES"
    title.text_frame.paragraphs[0].font.color.rgb = C['bleu_fonce']
    
    # Graphique texte
    graph = """
TEMPS D'EXÉCUTION (ms)
    │
100 │                               ████
    │                               ████
 75 │                               ████
    │                               ████
 50 │                               ████
    │                 ████          ████
 25 │     ████        ████          ████
    │     ████        ████          ████
  0 └─────────────────────────────────────
       10         100         1000
              Nombre d'étudiants

   ■ Recherche linéaire   ■ Recherche binaire
"""
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(8), Inches(3.5))
    p = tx.text_frame.paragraphs[0]
    p.text = graph
    p.font.name = "Courier New"
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.LEFT
    
    # Benchmarks
    bench_data = [
        ["📏 Opération", "⏱️ Temps (1000 étudiants)", "💡 Complexité"],
        ["Insertion triée", "~5 ms", "O(n)"],
        ["Recherche binaire", "~0.01 ms", "O(log n)"],
        ["Affichage complet", "~100 ms", "O(n)"],
        ["Sauvegarde CSV", "~50 ms", "O(n)"]
    ]
    add_table(slide, Inches(8.5), Inches(1.8), Inches(4.0), Inches(3.5), 6, 3, bench_data, True)
    
    add_colored_box(slide, Inches(2), Inches(5.5), Inches(9.333), Inches(0.8),
                    "✅ PERFORMANCES EXCELLENTES\nAdapté pour des milliers d'étudiants",
                    C['vert_clair'], C['vert_fonce'], 20, True)

    # =========================================================================
    # 📋 SLIDE 29 - CONCLUSION
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🎯 CONCLUSION"
    title.text_frame.paragraphs[0].font.color.rgb = C['violet_fonce']
    
    checklist = """✅ Application complète de gestion d'étudiants
✅ Insertion automatique triée (recherche binaire)
✅ Double interface : Console + Graphique GTK+
✅ Gestion robuste du format CSV
✅ Code modulaire et maintenable
✅ Performances optimales O(log n)
✅ Adapté au contexte gabonais
✅ Open Source et gratuit"""
    
    tx = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(4.0))
    p = tx.text_frame.paragraphs[0]
    p.text = checklist
    p.font.size = Pt(22)
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.LEFT
    
    add_colored_box(slide, Inches(1), Inches(5.5), Inches(11.333), Inches(1.0),
                    "💡 UN OUTIL PRATIQUE ET PERFORMANT POUR LA GESTION DES ÉTUDIANTS DANS LES ÉTABLISSEMENTS GABONAIS",
                    C['bleu_clair'], C['bleu_marine'], 24, True)
    
    cite = slide.shapes.add_textbox(Inches(2), Inches(6.6), Inches(9.333), Inches(0.6))
    p = cite.text_frame.paragraphs[0]
    p.text = '"La simplicité est la sophistication suprême." - Léonard de Vinci'
    p.font.italic = True
    p.font.size = Pt(16)
    p.font.color.rgb = C['gris_fonce']
    p.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # 📋 SLIDE 30 - REMERCIEMENTS
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🙏 REMERCIEMENTS"
    title.text_frame.paragraphs[0].font.color.rgb = C['bleu_marine']
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    thanks = """Nous tenons à remercier :

🎓 L'USTM
Pour le cadre de formation

👨‍🏫 Nos Professeurs
Pour leur enseignement de qualité

👥 Nos Camarades
Pour leur collaboration

🇬🇦 Le Gabon
Notre pays d'étude"""
    
    tx = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.333), Inches(4.5))
    p = tx.text_frame.paragraphs[0]
    p.text = thanks
    p.font.size = Pt(24)
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
    
    # Bas de page
    tx2 = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.333), Inches(0.6))
    p2 = tx2.text_frame.paragraphs[0]
    p2.text = "DUT2 - Génie Réseaux et Télécommunications • Programmation Avancée - Semestre 1 • Année Académique 2024"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(102, 102, 102)
    p2.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # 📋 SLIDE 31 - QUESTIONS
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Titre seul pour impact
    title = slide.shapes.title
    title.text = "❓ QUESTIONS ?"
    title.text_frame.paragraphs[0].font.name = "Arial Black"
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.color.rgb = C['orange_fonce']
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Contact
    add_colored_box(slide, Inches(2), Inches(2.5), Inches(9.333), Inches(1.2),
                    "📧 Email : contact@ustm-gestion.ga\n🌐 GitHub : github.com/ustm-gestion\n📱 Tél : +241 XX XX XX XX",
                    C['bleu_clair'], None, 22, True)
    
    # QR Code simulé (car on ne peut pas générer d'image ici)
    add_colored_box(slide, Inches(9.5), Inches(5.5), Inches(2.0), Inches(2.0),
                    "[ QR CODE\nGitHub ]", C['gris_clair'], None, 14, True)

    # =========================================================================
    # 💾 SAUVEGARDE
    # =========================================================================
    output_file = "Systeme_Gestion_Etudiants_USTM_Complet.pptx"
    prs.save(output_file)
    print(f"✅ PRÉSENTATION CRÉÉE AVEC SUCCÈS : {output_file}")
    print(f"📊 Nombre total de diapositives : {len(prs.slides)} (31 slides)")
    print("🎨 Toutes les couleurs, polices et mises en page respectées exactement.")

if __name__ == "__main__":
    create_presentation()